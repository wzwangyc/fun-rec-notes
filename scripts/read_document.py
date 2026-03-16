#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
办公文档内容读取工具
支持：Word, Excel, PPT, PDF, TXT, 图片 OCR
"""

import sys
import os
from pathlib import Path

def read_word(file_path):
    """读取 Word 文档 (.docx)"""
    import docx
    doc = docx.Document(file_path)
    content = []
    for para in doc.paragraphs:
        if para.text.strip():
            content.append(para.text)
    return "\n".join(content)

def read_excel(file_path, sheet=None):
    """读取 Excel 文件 (.xlsx, .xls)"""
    import openpyxl
    wb = openpyxl.load_workbook(file_path)
    if sheet:
        ws = wb[sheet]
    else:
        ws = wb.active
    content = []
    for row in ws.iter_rows(values_only=True):
        if any(cell is not None for cell in row):
            content.append("\t".join(str(cell) if cell else "" for cell in row))
    return "\n".join(content)

def read_ppt(file_path):
    """读取 PPT 文件 (.pptx)"""
    from pptx import Presentation
    prs = Presentation(file_path)
    content = []
    for i, slide in enumerate(prs.slides, 1):
        content.append(f"=== 幻灯片 {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text)
    return "\n".join(content)

def read_pdf(file_path):
    """读取 PDF 文件"""
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    content = []
    for i, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            content.append(f"=== 第 {i} 页 ===\n{text}")
    doc.close()
    return "\n".join(content)

def read_txt(file_path):
    """读取 TXT 文件"""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def read_image(file_path):
    """图片 OCR 文字识别"""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang='chi_sim+eng')
        return text
    except Exception as e:
        return f"OCR 识别失败：{e}"

def read_file(file_path):
    """根据文件扩展名自动选择读取方式"""
    path = Path(file_path)
    if not path.exists():
        return f"错误：文件不存在 - {file_path}"
    
    ext = path.suffix.lower()
    read_func = {
        '.docx': read_word,
        '.doc': read_word,
        '.xlsx': read_excel,
        '.xls': read_excel,
        '.pptx': read_ppt,
        '.ppt': read_ppt,
        '.pdf': read_pdf,
        '.txt': read_txt,
        '.md': read_txt,
        '.png': read_image,
        '.jpg': read_image,
        '.jpeg': read_image,
        '.gif': read_image,
        '.bmp': read_image,
    }.get(ext)
    
    if not read_func:
        return f"不支持的文件类型：{ext}"
    
    try:
        content = read_func(file_path)
        return f"文件：{path.name}\n类型：{ext}\n\n{content}"
    except Exception as e:
        return f"读取失败：{e}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法：python read_document.py <文件路径>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = read_file(file_path)
    print(result)
